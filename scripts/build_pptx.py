# -*- coding: utf-8 -*-
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
import copy
from PIL import Image

BASE_DIR = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
IMG_DIR = os.path.join(BASE_DIR, "パワポ用画像")
PROFILE_DIR = os.path.join(BASE_DIR, "自己紹介画像")

# ---- サイトのデザイントークン（public/style.css 準拠） ----
MAIN = RGBColor(0xB5, 0x55, 0x1F)       # --main-color テラコッタ
MAIN_HOVER = RGBColor(0x96, 0x46, 0x1A) # --main-color-hover
BG = RGBColor(0xF6, 0xF1, 0xE6)         # --bg
SURFACE = RGBColor(0xFF, 0xFD, 0xF8)    # --surface
SURFACE2 = RGBColor(0xEF, 0xE6, 0xD3)   # --surface-2
BORDER = RGBColor(0xDD, 0xCC, 0xAC)     # --border
INK = RGBColor(0x2C, 0x24, 0x18)        # --ink
SUB = RGBColor(0x7A, 0x6D, 0x54)        # --sub
STAMP = RGBColor(0x8A, 0x73, 0x55)      # --stamp
CREAM = RGBColor(0xFF, 0xF8, 0xEC)      # main上の文字色
TAGBG = RGBColor(0xF7, 0xE9, 0xDC)
TAGBORDER = RGBColor(0xEC, 0xD3, 0xBA)

FONT_SERIF = "Hiragino Mincho ProN"
FONT_SANS = "Hiragino Kaku Gothic ProN"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
BLANK = prs.slide_layouts[6]


def add_slide():
    return prs.slides.add_slide(BLANK)


def set_bg(slide, color=BG):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def no_line(shape):
    shape.line.fill.background()


def rect(slide, l, t, w, h, fill=None, line=None, line_w=None, shadow=False, radius=None):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    sp = slide.shapes.add_shape(shape_type, l, t, w, h)
    if radius:
        try:
            sp.adjustments[0] = radius
        except Exception:
            pass
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid()
        sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = line_w or Pt(1)
    sp.shadow.inherit = False
    return sp


def pill(slide, l, t, w, h, fill=None, line=None, line_w=None):
    sp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    try:
        sp.adjustments[0] = 0.5
    except Exception:
        pass
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid()
        sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = line_w or Pt(1)
    sp.shadow.inherit = False
    return sp


def textbox(slide, l, t, w, h, text, size=14, color=INK, bold=False, align=PP_ALIGN.LEFT,
            font=FONT_SANS, anchor=MSO_ANCHOR.TOP, line_spacing=1.0, wrap=True, italic=False,
            letter_spacing=None):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    lines = text.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        run = p.add_run()
        run.text = line
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = color
        run.font.name = font
    return tb


def multi_run_para(tf, first, runs_specs, align=PP_ALIGN.LEFT, line_spacing=1.0):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align
    p.line_spacing = line_spacing
    for txt, size, color, bold, font in runs_specs:
        run = p.add_run()
        run.text = txt
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
        run.font.name = font
    return p


def footer(slide, page_no, total=15, label="Webプログラミング体験授業"):
    textbox(slide, Inches(0.5), Inches(7.12), Inches(6), Inches(0.3),
            label, size=9, color=SUB, font=FONT_SANS)
    textbox(slide, Inches(12.0), Inches(7.12), Inches(0.9), Inches(0.3),
            f"{page_no} / {total}", size=9, color=SUB, font=FONT_SANS, align=PP_ALIGN.RIGHT)


def kicker(slide, text):
    """左上の小さい見出しラベル（例：STEP 1）"""
    bar = rect(slide, Inches(0.5), Inches(0.42), Inches(0.06), Inches(0.28), fill=MAIN)
    textbox(slide, Inches(0.66), Inches(0.4), Inches(4), Inches(0.32),
            text, size=13, color=MAIN, bold=True, font=FONT_SANS, letter_spacing=True)


def title(slide, text, size=28):
    textbox(slide, Inches(0.5), Inches(0.72), Inches(11.5), Inches(0.8),
            text, size=size, color=INK, bold=True, font=FONT_SERIF)


def brand_dot(slide, l, t, d=Inches(0.3)):
    dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, l, t, d, d)
    dot.fill.solid()
    dot.fill.fore_color.rgb = MAIN
    dot.line.fill.background()
    dot.shadow.inherit = False
    tf = dot.text_frame
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "旅"
    run.font.size = Pt(12)
    run.font.bold = True
    run.font.color.rgb = CREAM
    run.font.name = FONT_SANS
    return dot


def card(slide, l, t, w, h, fill=SURFACE, line=BORDER, radius=0.05):
    c = rect(slide, l, t, w, h, fill=fill, line=line, line_w=Pt(1), radius=radius)
    return c


def step_badge(slide, l, t, num, label):
    """円形バッジ＋ラベル"""
    d = Inches(0.42)
    dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, l, t, d, d)
    dot.fill.solid()
    dot.fill.fore_color.rgb = MAIN
    dot.line.fill.background()
    dot.shadow.inherit = False
    tf = dot.text_frame
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = str(num)
    run.font.size = Pt(15)
    run.font.bold = True
    run.font.color.rgb = CREAM
    run.font.name = FONT_SANS


def code_chip(slide, l, t, w, h, text):
    """ファイル名だけを表示するチップ"""
    c = rect(slide, l, t, w, h, fill=INK, radius=0.15)
    textbox(slide, l, t, w, h, text, size=13, color=RGBColor(0xF7, 0xE9, 0xDC), bold=False,
            align=PP_ALIGN.CENTER, font="Courier New", anchor=MSO_ANCHOR.MIDDLE)
    return c


CODE_DIM = RGBColor(0x9A, 0x9A, 0x8A)
CODE_TEXT = RGBColor(0xF7, 0xE9, 0xDC)
GUTTER_BG = RGBColor(0x14, 0x11, 0x0C)
GUTTER_TEXT = RGBColor(0xFF, 0xFF, 0xFF)


def code_block(slide, l, t, w, h, lines, code_size=12.5, gutter_w=Inches(0.55)):
    """行番号ガター付きのコードブロック。
    lines: [(行番号 or None, テキスト, 強調するか)] のリスト。行番号Noneなら空欄行として詰める。
    """
    rect(slide, l, t, w, h, fill=INK, radius=0.06)
    rect(slide, l, t, gutter_w, h, fill=GUTTER_BG, radius=0)
    # 角の丸みを合わせるため、ガター右端の直線部分だけ再度重ねて角を隠す簡易対応は不要
    # (rectは矩形なのでガターは左側面のみでOK。radius=0で四角のまま重ねる)
    line_h = Emu(int(h / len(lines)))
    for i, (num, text, emphasize) in enumerate(lines):
        y = Emu(t + line_h * i)
        if num is not None:
            textbox(slide, l, y, Emu(gutter_w - Inches(0.12)), line_h, str(num), size=code_size - 1,
                    color=GUTTER_TEXT, align=PP_ALIGN.RIGHT, font="Courier New", anchor=MSO_ANCHOR.MIDDLE)
        textbox(slide, Emu(l + gutter_w + Inches(0.18)), y, Emu(w - gutter_w - Inches(0.3)), line_h,
                text, size=code_size, color=(CODE_TEXT if emphasize else CODE_DIM),
                font="Courier New", anchor=MSO_ANCHOR.MIDDLE)


def arrow(slide, l, t, w, h, fill=STAMP):
    sp = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, l, t, w, h)
    sp.fill.solid()
    sp.fill.fore_color.rgb = fill
    sp.line.fill.background()
    sp.shadow.inherit = False
    return sp


def tag_badge(slide, l, t, w, h, text):
    c = pill(slide, l, t, w, h, fill=TAGBG, line=TAGBORDER, line_w=Pt(0.75))
    textbox(slide, l, t, w, h, text, size=12, color=MAIN, bold=False, align=PP_ALIGN.CENTER,
            font=FONT_SANS, anchor=MSO_ANCHOR.MIDDLE)
    return c


def screenshot(slide, path, l, t, max_w, max_h, caption=None, label=None, label_color=SUB, border=BORDER, frame=True, pad=Pt(4)):
    """スクリーンショットを枠付きカードで配置し、任意でラベル・キャプションを添える"""
    im = Image.open(path)
    iw, ih = im.size
    ratio = min(max_w / iw, max_h / ih)
    w = Emu(int(iw * ratio))
    h = Emu(int(ih * ratio))
    x = Emu(int(l + (max_w - w) / 2))
    y = t
    if label:
        label_h = Inches(0.3)
        label_gap = Inches(0.06)
        label_bottom = Emu(y - pad - label_gap)
        textbox(slide, l, Emu(label_bottom - label_h), max_w, label_h, label, size=11.5,
                color=label_color, bold=True, font=FONT_SANS, anchor=MSO_ANCHOR.BOTTOM)
    if frame:
        rect(slide, Emu(x - pad), Emu(y - pad), Emu(w + pad * 2), Emu(h + pad * 2),
             fill=SURFACE, line=border, line_w=Pt(1), radius=0.04)
    slide.shapes.add_picture(path, x, y, width=w, height=h)
    bottom = Emu(y + h)
    if caption:
        textbox(slide, l, Emu(bottom + Inches(0.12)), max_w, Inches(0.35), caption, size=11,
                color=SUB, align=PP_ALIGN.CENTER, font=FONT_SANS)
    return bottom


def bullet_block(slide, l, t, w, h, lines, size=13.5, color=INK, gap=0.42, bold_first=False, marker="●"):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = 1.25
        p.space_after = Pt(10)
        run = p.add_run()
        run.text = f"{marker}  {line}"
        run.font.size = Pt(size)
        run.font.color.rgb = color
        run.font.name = FONT_SANS
    return tb


# ============================================================
# スライド 1: タイトル
# ============================================================
s = add_slide()
set_bg(s, BG)
# 装飾: 右下・左上に淡い円
c1 = slide_deco = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(-2), Inches(-2), Inches(5), Inches(5))
c1.fill.solid(); c1.fill.fore_color.rgb = SURFACE2; c1.fill.transparency = 0
c1.line.fill.background(); c1.shadow.inherit = False
# 半透明感を出すため薄い色に留める
c1.fill.fore_color.rgb = RGBColor(0xEF, 0xE6, 0xD3)

c2 = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(10.8), Inches(4.8), Inches(4.2), Inches(4.2))
c2.fill.solid(); c2.fill.fore_color.rgb = RGBColor(0xF7, 0xE9, 0xDC)
c2.line.fill.background(); c2.shadow.inherit = False

brand_dot(s, Inches(0.9), Inches(0.9), Inches(0.5))
textbox(s, Inches(1.55), Inches(0.9), Inches(4), Inches(0.5), "旅ノート | Web体験版",
        size=15, color=SUB, bold=True, font=FONT_SANS, anchor=MSO_ANCHOR.MIDDLE)

textbox(s, Inches(0.9), Inches(2.5), Inches(11.5), Inches(1.3),
        "Webプログラミング体験授業", size=42, color=INK, bold=True, font=FONT_SERIF)
textbox(s, Inches(0.9), Inches(3.5), Inches(11.5), Inches(1.0),
        "「旅行スポット共有サイト」を、少しずつ完成させよう", size=20, color=MAIN, bold=True, font=FONT_SERIF)

line = rect(s, Inches(0.9), Inches(4.35), Inches(2.4), Pt(3), fill=MAIN)

textbox(s, Inches(0.9), Inches(4.65), Inches(9), Inches(0.6),
        "コードを直すと、Webアプリの動きが変わる。それを自分の手で体感しよう。",
        size=15, color=SUB, font=FONT_SANS)

# 下部：使用ツールのタグ
tools = ["VSCode", "PHP", "SQLite", "HTML / CSS / JS"]
x = Inches(0.9)
for t_ in tools:
    w = Inches(0.5 + len(t_) * 0.11)
    tag_badge(s, x, Inches(6.3), w, Inches(0.42), t_)
    x = Emu(x + w + Inches(0.18))

footer(s, 1, total=21)

# ============================================================
# スライド 1.5: 自己紹介（曽根先生・2026年版）
# ============================================================
s = add_slide()
set_bg(s)
kicker(s, "自己紹介")
title(s, "曽根　大智　(Sone Daichi)")

img_path = os.path.join(PROFILE_DIR, "曽根先生_2026_実習授業.png")
screenshot(s, img_path, Inches(7.55), Inches(1.9), Inches(4.25), Inches(4.6),
           caption="高等専修学校勤務時：実習授業での一コマ")

card(s, Inches(0.5), Inches(1.9), Inches(6.7), Inches(4.6))
textbox(s, Inches(0.85), Inches(2.15), Inches(6.0), Inches(0.4),
        "大阪情報コンピュータ専門学校 講師", size=14, bold=True, color=MAIN, font=FONT_SANS)
bullet_block(s, Inches(0.85), Inches(2.7), Inches(6.05), Inches(3.5),
             ["今年3月まで大阪情報コンピュータ高等専修学校\n（OICのグループ校）に勤務",
              "高校生に情報系の授業を教えていた",
              "実は高等専修学校のパンフレットにも\n実習授業の一コマがちらっと掲載されている"],
             size=14.5, gap=0.5)

footer(s, 2, total=21)

# ============================================================
# スライド 1.6: 自己紹介（厨子先生・2024年版・予備）
# ============================================================
s = add_slide()
set_bg(s)
kicker(s, "自己紹介（予備）")
title(s, "厨子　直人　(Zushi, Naoto)")

img_path = os.path.join(PROFILE_DIR, "厨子先生_2024_杖を持った写真_conv.jpg")
screenshot(s, img_path, Inches(8.35), Inches(1.9), Inches(3.45), Inches(4.6))

card(s, Inches(0.5), Inches(1.9), Inches(7.5), Inches(2.55))
textbox(s, Inches(0.85), Inches(2.12), Inches(6.9), Inches(0.4),
        "経歴", size=14, bold=True, color=MAIN, font=FONT_SANS)
bullet_block(s, Inches(0.85), Inches(2.6), Inches(6.9), Inches(1.75),
             ["(株)Sky、大阪府四條畷市役所などに所属。NTT西日本などで\nプロジェクトマネジメント・情報セキュリティ・アーキテクトとして従事",
              "日本で初めての予測変換機能を実装した経歴を持つ",
              "コンピュータ専門誌（「I/O」など）で記事を執筆"],
             size=13)

card(s, Inches(0.5), Inches(4.65), Inches(7.5), Inches(1.85))
textbox(s, Inches(0.85), Inches(4.87), Inches(6.9), Inches(0.4),
        "本校での担当授業", size=14, bold=True, color=MAIN, font=FONT_SANS)
bullet_block(s, Inches(0.85), Inches(5.3), Inches(6.9), Inches(1.15),
             ["IT基礎科目／プログラミング（1年生）、システム設計各種（1〜3年生）",
              "基本情報／応用情報対策講座、AIに関する授業、クラス担任"],
             size=13)

footer(s, 3, total=21)

# ============================================================
# スライド 2: 今日やること（アジェンダ）
# ============================================================
s = add_slide()
set_bg(s)
kicker(s, "AGENDA")
title(s, "今日やること")
textbox(s, Inches(0.5), Inches(1.35), Inches(10.5), Inches(0.5),
        "旅行スポット共有サイトの「未完成な部分」を、順番に完成させていきます",
        size=14, color=SUB, font=FONT_SANS)

items = [
    ("0", "プログラミングって何？", "料理づくりに例えて、システム開発の流れをつかむ"),
    ("1", "コードを直してみよう", "ボタンの文字を変える／コメントアウトを外す"),
    ("2", "APIってなに？", "現在地から住所が自動で入力される仕組みを体験"),
    ("3", "データベースって何？", "投稿やコメントの情報がどう保存されているか"),
    ("4", "入力チェックを作ろう", "空欄のまま投稿できてしまう問題を自分で直す"),
    ("5", "自由にアレンジ", "表示の文言などを自分の工夫で変えてみる"),
]

col_w = Inches(3.9)
row_h = Inches(1.55)
grid_w = Emu(col_w * 3 + Inches(0.18) * 2)
start_x = Emu((SLIDE_W - grid_w) // 2)
start_y = Inches(2.05)
gap_x = Inches(0.18)
gap_y = Inches(0.2)

for i, (num, ttl, desc) in enumerate(items):
    row = i // 3
    col = i % 3
    x = Emu(start_x + col * (col_w + gap_x))
    y = Emu(start_y + row * (row_h + gap_y))
    card(s, x, y, col_w, row_h)
    step_badge(s, Emu(x + Inches(0.22)), Emu(y + Inches(0.2)), num, "")
    textbox(s, Emu(x + Inches(0.8)), Emu(y + Inches(0.18)), Emu(col_w - Inches(1.0)), Inches(0.5),
            ttl, size=14.5, color=INK, bold=True, font=FONT_SANS)
    textbox(s, Emu(x + Inches(0.22)), Emu(y + Inches(0.78)), Emu(col_w - Inches(0.4)), Inches(0.7),
            desc, size=11.5, color=SUB, font=FONT_SANS, line_spacing=1.25)

footer(s, 4, total=21)

# ============================================================
# スライド 3: プログラミングって何？（料理の比喩）
# ============================================================
s = add_slide()
set_bg(s)
kicker(s, "STEP 0")
title(s, "プログラミングって何？")
textbox(s, Inches(0.5), Inches(1.35), Inches(11), Inches(0.5),
        "システムをつくる流れは、「料理をつくる」流れとよく似ています",
        size=14, color=SUB, font=FONT_SANS)

steps = [
    ("要件定義", "メニューを決める"),
    ("設計", "レシピを考える"),
    ("実装", "調理する"),
    ("テスト", "味見して直す"),
    ("運用保守", "食事・片付け"),
]
n = len(steps)
total_w = Inches(11.3)
box_w = Inches(1.95)
gap = Emu((total_w - box_w * n) // (n - 1))
x = Emu((SLIDE_W - total_w) // 2)
y = Inches(2.35)
box_h = Inches(1.5)

for i, (top_label, bottom_label) in enumerate(steps):
    card(s, x, y, box_w, box_h, fill=SURFACE if i != 2 else SURFACE2, line=BORDER if i != 2 else MAIN)
    textbox(s, x, Emu(y + Inches(0.18)), box_w, Inches(0.4), top_label, size=13.5, bold=True,
            color=INK if i != 2 else MAIN, align=PP_ALIGN.CENTER, font=FONT_SANS)
    rect(s, Emu(x + Inches(0.65)), Emu(y + Inches(0.62)), Inches(0.65), Pt(1.5), fill=BORDER)
    textbox(s, Emu(x + Inches(0.08)), Emu(y + Inches(0.82)), Emu(box_w - Inches(0.16)), Inches(0.6),
            bottom_label, size=11.5, color=SUB, align=PP_ALIGN.CENTER, font=FONT_SANS, line_spacing=1.15)
    if i < n - 1:
        arrow(s, Emu(x + box_w + Inches(0.04)), Emu(y + Inches(0.6)), Emu(gap - Inches(0.08)), Inches(0.3))
    x = Emu(x + box_w + gap)

badge = pill(s, Inches(3.85), Inches(4.35), Inches(5.6), Inches(0.55), fill=MAIN)
textbox(s, Inches(3.85), Inches(4.35), Inches(5.6), Inches(0.55),
        "今日は「実装」「テスト」の部分を体験します", size=14.5, color=CREAM, bold=True,
        align=PP_ALIGN.CENTER, font=FONT_SANS, anchor=MSO_ANCHOR.MIDDLE)

footer(s, 5, total=21)

# ============================================================
# スライド 4: 表の顔・裏の顔（ラーメン比喩）
# ============================================================
s = add_slide()
set_bg(s)
kicker(s, "STEP 0")
title(s, "Webページの「表の顔」と「裏の顔」")
textbox(s, Inches(0.5), Inches(1.35), Inches(11), Inches(0.5),
        "見えている部分（フロントエンド）と、見えない部分（バックエンド）があります",
        size=14, color=SUB, font=FONT_SANS)

# 左カード：カップラーメン
lx, ly, lw, lh = Inches(1.0), Inches(2.1), Inches(5.55), Inches(4.5)
card(s, lx, ly, lw, lh)
rect(s, lx, ly, lw, Inches(0.1), fill=STAMP)
textbox(s, Emu(lx + Inches(0.35)), Emu(ly + Inches(0.3)), Emu(lw - Inches(0.7)), Inches(0.45),
        "🍜 表の顔だけ＝カップラーメン", size=17, color=INK, bold=True, font=FONT_SANS)
textbox(s, Emu(lx + Inches(0.35)), Emu(ly + Inches(0.85)), Emu(lw - Inches(0.7)), Inches(0.4),
        "フロントエンド（HTML / CSS / JavaScript）のみ", size=12.5, color=SUB, font=FONT_SANS)
bullet_block(s, Emu(lx + Inches(0.35)), Emu(ly + Inches(1.5)), Emu(lw - Inches(0.7)), Inches(2.6),
             ["すぐに提供できる", "だれに対しても同じ表示・同じ動き", "サーバー側の処理がない＝\n個人に合わせた表示はできない"],
             size=13)

# 右カード：ラーメン屋
rx, ry, rw, rh = Inches(6.78), Inches(2.1), Inches(5.55), Inches(4.5)
card(s, rx, ry, rw, rh, fill=RGBColor(0xF7, 0xE9, 0xDC), line=MAIN)
rect(s, rx, ry, rw, Inches(0.1), fill=MAIN)
textbox(s, Emu(rx + Inches(0.35)), Emu(ry + Inches(0.3)), Emu(rw - Inches(0.7)), Inches(0.45),
        "🍜 裏の顔まで＝ラーメン屋さん", size=17, color=INK, bold=True, font=FONT_SANS)
textbox(s, Emu(rx + Inches(0.35)), Emu(ry + Inches(0.85)), Emu(rw - Inches(0.7)), Inches(0.4),
        "＋ バックエンド（今回は PHP）", size=12.5, color=MAIN, bold=True, font=FONT_SANS)
bullet_block(s, Emu(rx + Inches(0.35)), Emu(ry + Inches(1.5)), Emu(rw - Inches(0.7)), Inches(2.6),
             ["提供までに少し時間がかかる", "お客さんに合わせた味・トッピングができる", "サーバー側でデータを処理\n＝一人ひとりに合わせた表示ができる"],
             size=13)

textbox(s, Inches(0.5), Inches(6.75), Inches(11.3), Inches(0.4),
        "今日は、この「厨房エリア（バックエンド）」を覗いてみましょう！",
        size=13.5, color=MAIN, bold=True, align=PP_ALIGN.CENTER, font=FONT_SANS)

footer(s, 6, total=21)

# ============================================================
# スライド 5: 目の前の画面を見てみよう
# ============================================================
s = add_slide()
set_bg(s)
kicker(s, "STEP 0 → 1")
title(s, "目の前の画面を見てみよう")
textbox(s, Inches(0.5), Inches(1.35), Inches(11), Inches(0.5),
        "VSCodeとブラウザは、すでに開いた状態からスタートします",
        size=14, color=SUB, font=FONT_SANS)

bullet_block(s, Inches(0.85), Inches(2.35), Inches(5.4), Inches(3.2),
             ["ブラウザに表示されているのが「キミの旅」のトップ画面",
              "VSCodeにはプロジェクトのファイル一式がすでに開いている",
              "エディタでコードを直して保存 → ブラウザを再読み込み、\nの繰り返しで進めていく"],
             size=14, gap=0.5)

rx, ry, rw, rh = Inches(6.55), Inches(2.15), Inches(5.25), Inches(4.4)
screenshot(s, os.path.join(IMG_DIR, "01_トップ画面.png"), rx, ry, rw, rh,
           caption="トップ画面（一覧・検索）")

footer(s, 7, total=21)

# ============================================================
# スライド 6: ボタンの文字を変えてみよう
# ============================================================
s = add_slide()
set_bg(s)
kicker(s, "STEP 1")
title(s, "ボタンの文字を変えてみよう")
textbox(s, Inches(0.5), Inches(1.32), Inches(7.5), Inches(0.4),
        "投稿ボタンの表示が分かりにくい…",
        size=13.5, color=SUB, font=FONT_SANS)
code_chip(s, Inches(9.5), Inches(1.28), Inches(2.3), Inches(0.42), "login.php")

code_block(s, Inches(0.5), Inches(1.85), Inches(11.3), Inches(0.5),
           [(53, '<button ... class="btn-ghost">ボタン</button>', True)])

# Before/After（郵便番号欄〜ボタンにズームしたスクリーンショット）
screenshot(s, os.path.join(IMG_DIR, "02_新規登録_修正前_ボタン文言_ズーム.png"), Inches(0.9), Inches(2.95),
           Inches(11.5), Inches(1.6), label="BEFORE：「ボタン」のまま", label_color=SUB)

screenshot(s, os.path.join(IMG_DIR, "02_新規登録_修正後_住所を検索_ズーム.png"), Inches(0.9), Inches(5.25),
           Inches(11.5), Inches(1.6), label="AFTER：「住所を検索」に変更", label_color=MAIN, border=MAIN)

footer(s, 8, total=21)

# ============================================================
# スライド 7: コメントアウトを外して住所検索を有効化
# ============================================================
s = add_slide()
set_bg(s)
kicker(s, "STEP 1")
title(s, "コメントアウトを外してみよう")
textbox(s, Inches(0.5), Inches(1.35), Inches(11), Inches(0.5),
        "ボタンを押しても、まだ住所は出てこない",
        size=14, color=SUB, font=FONT_SANS)

code_chip(s, Inches(0.5), Inches(2.0), Inches(3.4), Inches(0.5), "postal_lookup.php")
textbox(s, Inches(4.1), Inches(2.0), Inches(2.7), Inches(0.5),
        "/* */ を外そう", size=14.5, color=INK, bold=True, font=FONT_SANS,
        anchor=MSO_ANCHOR.MIDDLE)

code_block(s, Inches(0.5), Inches(2.65), Inches(6.15), Inches(1.65), [
    (25, "/*  ← ここから", False),
    (26, "find_postal_address($pdo, ...);", True),
    (40, "*/  ← ここまでを無効化している", False),
])

textbox(s, Inches(0.5), Inches(4.65), Inches(6.15), Inches(1.1),
        "「/* と */ で囲まれた部分は実行時に無視される」\nこの2文字を外すだけで、住所検索の関数呼び出しが有効になる。",
        size=13, color=SUB, font=FONT_SANS, line_spacing=1.35)

screenshot(s, os.path.join(IMG_DIR, "07_郵便番号検索_実行結果_ズーム.png"), Inches(6.95), Inches(2.3),
           Inches(5.85), Inches(4.65), label="有効化後：郵便番号から住所を自動入力", label_color=MAIN, border=MAIN)

footer(s, 9, total=21)

# ============================================================
# スライド 7.5: 裏側で起きていること（関数の考え方）
# ============================================================
s = add_slide()
set_bg(s)
kicker(s, "STEP 1")
title(s, "裏側で起きていること：関数の考え方")
textbox(s, Inches(0.5), Inches(1.35), Inches(11), Inches(0.5),
        "「郵便番号を渡すと、住所が返ってくる」——これが「関数」です",
        size=14, color=SUB, font=FONT_SANS)

cy = Inches(2.6)
ch = Inches(1.6)
# システム（呼び出す側）
card(s, Inches(0.5), cy, Inches(3.0), ch, fill=SURFACE2)
textbox(s, Inches(0.5), Emu(cy + Inches(0.15)), Inches(3.0), Inches(0.4), "あなた", size=14, bold=True,
        color=INK, align=PP_ALIGN.CENTER, font=FONT_SANS)
textbox(s, Inches(0.5), Emu(cy + Inches(0.6)), Inches(3.0), Inches(0.9), "「住所教えて」\nと呼び出す", size=12.5,
        color=SUB, align=PP_ALIGN.CENTER, font=FONT_SANS, line_spacing=1.2)

arrow(s, Inches(3.65), Emu(cy + Inches(0.15)), Inches(1.15), Inches(0.4))
textbox(s, Inches(3.55), Emu(cy - Inches(0.32)), Inches(1.4), Inches(0.35), "郵便番号を渡す",
        size=10.5, color=SUB, align=PP_ALIGN.CENTER, font=FONT_SANS)

# 関数（住所検索）
fx = Inches(4.95)
card(s, fx, cy, Inches(3.35), ch, fill=RGBColor(0xF7, 0xE9, 0xDC), line=MAIN)
textbox(s, fx, Emu(cy + Inches(0.15)), Inches(3.35), Inches(0.4), "find_postal_address()", size=13, bold=True,
        color=MAIN, align=PP_ALIGN.CENTER, font="Courier New")
textbox(s, fx, Emu(cy + Inches(0.6)), Inches(3.35), Inches(0.9), "（関数）\n中でどう処理するかは\n呼び出す側は知らなくてよい",
        size=11.5, color=SUB, align=PP_ALIGN.CENTER, font=FONT_SANS, line_spacing=1.2)

arrow(s, Inches(8.45), Emu(cy + Inches(0.15)), Inches(1.15), Inches(0.4))
textbox(s, Inches(8.3), Emu(cy - Inches(0.32)), Inches(1.5), Inches(0.35), "住所データを返す",
        size=10.5, color=SUB, align=PP_ALIGN.CENTER, font=FONT_SANS)

# 結果
card(s, Inches(9.75), cy, Inches(3.0), ch, fill=SURFACE2)
textbox(s, Inches(9.75), Emu(cy + Inches(0.15)), Inches(3.0), Inches(0.4), "画面", size=14, bold=True,
        color=INK, align=PP_ALIGN.CENTER, font=FONT_SANS)
textbox(s, Inches(9.75), Emu(cy + Inches(0.6)), Inches(3.0), Inches(0.9), "住所が\n自動入力される",
        size=12.5, color=SUB, align=PP_ALIGN.CENTER, font=FONT_SANS, line_spacing=1.2)

box = card(s, Inches(0.85), Inches(5.1), Inches(11.6), Inches(1.15), fill=SURFACE)
textbox(s, Inches(1.15), Inches(5.35), Inches(11.0), Inches(0.7),
        "「機能をひとつのまとまりとして作っておき、必要なときに呼び出す」——これが関数の考え方。\n機能をつくることも、プログラマーの大事な仕事の一つです。",
        13.5, color=INK, font=FONT_SANS, line_spacing=1.35)

footer(s, 10, total=21)

# ============================================================
# スライド 8: APIってなに？（現在地→住所）
# ============================================================
s = add_slide()
set_bg(s)
kicker(s, "STEP 2")
title(s, "APIってなに？")
textbox(s, Inches(0.5), Inches(1.35), Inches(11), Inches(0.5),
        "「📍 現在地を取得」ボタンを押すと、住所が自動で入力される仕組みを見てみよう",
        size=14, color=SUB, font=FONT_SANS)

# フロー図：4ステップ（コンパクト）
flow_items = [
    ("📍", "現在地(緯度・経度)を取得"),
    ("🖥", "geocode.phpが外部APIへ問い合わせ"),
    ("🌏", "外部の住所検索API\n(OpenStreetMap)"),
    ("📝", "住所が自動入力される"),
]
n = len(flow_items)
box_w = Inches(2.6)
total_w = Inches(11.3)
gap = Emu((total_w - box_w * n) // (n - 1))
x = Emu((SLIDE_W - total_w) // 2)
y = Inches(2.0)
box_h = Inches(1.0)

for i, (emoji, label) in enumerate(flow_items):
    fill = RGBColor(0xF7, 0xE9, 0xDC) if i == 2 else SURFACE
    line = MAIN if i == 2 else BORDER
    card(s, x, y, box_w, box_h, fill=fill, line=line)
    textbox(s, x, Emu(y + Inches(0.1)), Inches(0.65), Emu(box_h - Inches(0.2)), emoji, size=22,
            align=PP_ALIGN.CENTER, font=FONT_SANS, anchor=MSO_ANCHOR.MIDDLE)
    textbox(s, Emu(x + Inches(0.6)), y, Emu(box_w - Inches(0.7)), box_h,
            label, size=10.5, color=INK, align=PP_ALIGN.LEFT, font=FONT_SANS, line_spacing=1.15,
            anchor=MSO_ANCHOR.MIDDLE)
    if i < n - 1:
        arrow(s, Emu(x + box_w + Inches(0.04)), Emu(y + box_h/2 - Inches(0.12)), Emu(gap - Inches(0.08)), Inches(0.24))
    x = Emu(x + box_w + gap)

screenshot(s, os.path.join(IMG_DIR, "04_現在地取得_押す前_ズーム.png"), Inches(0.9), Inches(3.45),
           Inches(11.5), Inches(1.5), label="押す前", label_color=SUB)
screenshot(s, os.path.join(IMG_DIR, "04_現在地取得_押した後_ズーム.png"), Inches(0.9), Inches(5.45),
           Inches(11.5), Inches(1.5), label="押した後：住所が自動入力される", label_color=MAIN, border=MAIN)

footer(s, 11, total=21)

# ============================================================
# スライド 8.5: 郵便番号検索 vs 現在地取得（関数とAPIの比較）
# ============================================================
s = add_slide()
set_bg(s)
kicker(s, "STEP 2")
title(s, "さっきの関数と、何が違う？")
textbox(s, Inches(0.5), Inches(1.35), Inches(11), Inches(0.5),
        "どちらも「渡す→調べる→返ってくる」という形は同じ",
        size=14, color=SUB, font=FONT_SANS)

cy = Inches(2.5)
ch = Inches(2.1)

# 左：郵便番号検索（関数）
lx = Inches(0.85)
card(s, lx, cy, Inches(5.4), ch, fill=SURFACE2)
textbox(s, Emu(lx + Inches(0.3)), Emu(cy + Inches(0.22)), Inches(4.8), Inches(0.4),
        "郵便番号検索", size=15, bold=True, color=INK, font=FONT_SANS)
textbox(s, Emu(lx + Inches(0.3)), Emu(cy + Inches(0.68)), Inches(4.8), Inches(0.35),
        "find_postal_address()", size=12, color=SUB, font="Courier New")
bullet_block(s, Emu(lx + Inches(0.3)), Emu(cy + Inches(1.15)), Inches(4.8), Inches(0.9),
             ["調べる相手：自分のパソコンの中\n（postal_codes テーブル）"], size=12.5, gap=0.4)

# 右：現在地取得（外部API）
rx = Inches(6.95)
card(s, rx, cy, Inches(5.4), ch, fill=RGBColor(0xF7, 0xE9, 0xDC), line=MAIN)
textbox(s, Emu(rx + Inches(0.3)), Emu(cy + Inches(0.22)), Inches(4.8), Inches(0.4),
        "現在地取得（外部API）", size=15, bold=True, color=MAIN, font=FONT_SANS)
textbox(s, Emu(rx + Inches(0.3)), Emu(cy + Inches(0.68)), Inches(4.8), Inches(0.35),
        "OpenStreetMap Nominatim", size=12, color=MAIN, font="Courier New")
bullet_block(s, Emu(rx + Inches(0.3)), Emu(cy + Inches(1.15)), Inches(4.8), Inches(0.9),
             ["調べる相手：インターネットの向こうの\n別のサーバー"], size=12.5, gap=0.4)

box = card(s, Inches(0.85), Inches(4.95), Inches(11.6), Inches(1.4), fill=SURFACE)
textbox(s, Inches(1.15), Inches(5.15), Inches(11.0), Inches(1.0),
        "「値を渡す → 誰か（何か）が調べる → 答えが返ってくる」という形はまったく同じ。\n「API」という言葉はどちらの場所にも使えるが、実際に「APIを使う」というときは\nたいてい後者（＝Web API、インターネットの向こう）を指すことが多い。",
        13, color=INK, font=FONT_SANS, line_spacing=1.4)

footer(s, 12, total=21)

# ============================================================
# スライド 10: データベースって何？（comments.message を有効化）
# ============================================================
s = add_slide()
set_bg(s)
kicker(s, "STEP 3")
title(s, "データベースって何？")
textbox(s, Inches(0.5), Inches(1.35), Inches(11), Inches(0.5),
        "コメントは投稿できたのに、本文だけがどれも空欄になっている",
        size=14, color=SUB, font=FONT_SANS)

# Before/After（コメント表示、横並びで大きめ）
screenshot(s, os.path.join(IMG_DIR, "03_コメント欄_修正前_本文空欄.png"), Inches(0.9), Inches(2.35),
           Inches(5.35), Inches(2.5), label="BEFORE：本文が空欄", label_color=SUB)
screenshot(s, os.path.join(IMG_DIR, "03_コメント欄_修正後_本文表示.png"), Inches(7.05), Inches(2.35),
           Inches(5.35), Inches(2.5), label="AFTER：本文が表示される", label_color=MAIN, border=MAIN)

# 考えてみよう（問いかけのみ、答えは出さない）
card(s, Inches(0.85), Inches(5.4), Inches(11.6), Inches(1.15), fill=RGBColor(0xF7, 0xE9, 0xDC), line=MAIN)
textbox(s, Inches(1.2), Inches(5.65), Inches(10.9), Inches(0.7),
        "考えてみよう：「どのコメントか」を特定して正しく表示するには、どんな情報が必要そう？",
        size=15, bold=True, color=MAIN, font=FONT_SANS, anchor=MSO_ANCHOR.MIDDLE)

footer(s, 13, total=21)

# ============================================================
# スライド 10.5: データベースって何？（comments テーブルを覗く）
# ============================================================
s = add_slide()
set_bg(s)
kicker(s, "STEP 3")
title(s, "comments テーブルを覗いてみよう")
textbox(s, Inches(0.5), Inches(1.32), Inches(7.5), Inches(0.4),
        "実際のデータを見ながら、必要なカラムを確認する",
        size=13.5, color=SUB, font=FONT_SANS)
code_chip(s, Inches(9.5), Inches(1.28), Inches(2.3), Inches(0.42), "show.php")

screenshot(s, os.path.join(IMG_DIR, "08_DBテーブルの中身_comments_ズーム.png"), Inches(2.66), Inches(2.4),
           Inches(8.0), Inches(2.9), label="comments テーブルの中身（VSCode）", label_color=SUB, pad=Pt(16))

fields = ["コメントID", "投稿されたスポット", "投稿したユーザー", "コメント本文", "投稿日時"]
total_fw = sum(Inches(0.35 + len(f) * 0.16) for f in fields) + Inches(0.15) * (len(fields) - 1)
fx = Emu((SLIDE_W - total_fw) // 2)
fy = Inches(5.75)
for f in fields:
    w = Inches(0.35 + len(f) * 0.16)
    tag_badge(s, fx, fy, w, Inches(0.45), f)
    fx = Emu(fx + w + Inches(0.15))

# SQL文についてさらっと触れる
code_block(s, Inches(0.5), Inches(6.45), Inches(9.3), Inches(0.55), [
    (None, "SELECT * FROM comments WHERE spot_id = 12", True),
])
textbox(s, Inches(9.95), Inches(6.45), Inches(2.85), Inches(0.55),
        "この「SELECT〜」がSQL文\n（データベースへの問い合わせ文）",
        size=10.5, color=SUB, font=FONT_SANS, line_spacing=1.25, anchor=MSO_ANCHOR.MIDDLE)

footer(s, 14, total=21)

# ============================================================
# スライド 11: 入力チェックを追加しよう
# ============================================================
s = add_slide()
set_bg(s)
kicker(s, "STEP 4")
title(s, "入力チェックを追加しよう")
textbox(s, Inches(0.5), Inches(1.32), Inches(7.5), Inches(0.4),
        "タイトル欄が空白のままでも投稿できてしまう",
        size=13.5, color=SUB, font=FONT_SANS)
code_chip(s, Inches(9.5), Inches(1.28), Inches(2.3), Inches(0.42), "spots.php")

screenshot(s, os.path.join(IMG_DIR, "05_タイトル空欄カード_ズーム.png"), Inches(2.83), Inches(2.1),
           Inches(7.7), Inches(3.55), label="😮 左のカードだけタイトルが空欄（右は比較用）", label_color=MAIN, border=MAIN)

# 対応方針（画像の下に横並び）
lx = Inches(0.5)
card(s, lx, Inches(6.15), Inches(5.6), Inches(0.85))
textbox(s, Emu(lx + Inches(0.3)), Inches(6.28), Inches(5.0), Inches(0.35), "✅ やること", size=13.5,
        bold=True, color=INK, font=FONT_SANS)
bullet_block(s, Emu(lx + Inches(0.3)), Inches(6.6), Inches(5.0), Inches(0.35),
             ["if文（/* */ でコメントアウト）を有効化"], size=11, gap=0.4)

rx = Inches(6.3)
code_block(s, rx, Inches(6.15), Inches(6.5), Inches(0.85), [
    (32, "if ($title === '') {", True),
])

footer(s, 15, total=21)

# ============================================================
# スライド 11.5: 発展トーク（サニタイズ・SQLインジェクション）
# ============================================================
s = add_slide()
set_bg(s)
kicker(s, "発展トーク")
title(s, "HTMLは「タグ」でできている")
textbox(s, Inches(0.5), Inches(1.35), Inches(11.5), Inches(0.5),
        "<> で囲まれた「タグ」を組み合わせて、ブラウザに表示するページを作っている",
        size=15, color=SUB, font=FONT_SANS)

# コード → 表示結果の対比
lx = Inches(0.5)
card(s, lx, Inches(2.15), Inches(5.9), Inches(3.8))
textbox(s, Emu(lx + Inches(0.3)), Inches(2.35), Inches(5.3), Inches(0.4),
        "書いたコード", size=14, bold=True, color=MAIN, font=FONT_SANS)
codebox = rect(s, Emu(lx + Inches(0.3)), Inches(2.85), Inches(5.3), Inches(2.7), fill=INK, radius=0.06)
tf = codebox.text_frame
tf.word_wrap = True
tf.margin_left = Inches(0.28); tf.margin_top = Inches(0.25)
lines_code = [
    [("<h1>", RGBColor(0xF7, 0xE9, 0xDC)), ("  ← 見出しタグ", RGBColor(0x9A, 0x9A, 0x8A))],
    [("  地元の人しか知らない絶景カフェ", RGBColor(0xE0, 0xB0, 0x88))],
    [("</h1>", RGBColor(0xF7, 0xE9, 0xDC))],
    [("<p>", RGBColor(0xF7, 0xE9, 0xDC)), ("  ← 本文タグ", RGBColor(0x9A, 0x9A, 0x8A))],
    [("  夕方の光が綺麗です", RGBColor(0xE0, 0xB0, 0x88))],
    [("</p>", RGBColor(0xF7, 0xE9, 0xDC))],
]
for i, runs in enumerate(lines_code):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.line_spacing = 1.4
    p.alignment = PP_ALIGN.LEFT
    for txt, color in runs:
        r = p.add_run(); r.text = txt
        r.font.name = "Courier New"; r.font.size = Pt(13.5); r.font.color.rgb = color

rx = Inches(6.95)
card(s, rx, Inches(2.15), Inches(5.9), Inches(3.8), fill=RGBColor(0xF7, 0xE9, 0xDC), line=MAIN)
textbox(s, Emu(rx + Inches(0.3)), Inches(2.35), Inches(5.3), Inches(0.4),
        "ブラウザでの見え方", size=14, bold=True, color=MAIN, font=FONT_SANS)
disp = card(s, Emu(rx + Inches(0.3)), Inches(2.85), Inches(5.3), Inches(2.7), fill=SURFACE)
textbox(s, Emu(rx + Inches(0.55)), Inches(3.1), Inches(4.8), Inches(0.6),
        "地元の人しか知らない絶景カフェ", size=19, bold=True, color=INK, font=FONT_SERIF)
textbox(s, Emu(rx + Inches(0.55)), Inches(3.9), Inches(4.8), Inches(0.5),
        "夕方の光が綺麗です", size=14, color=SUB, font=FONT_SANS)

textbox(s, Inches(0.5), Inches(6.15), Inches(12.3), Inches(0.7),
        "見出しタグ・本文タグなど、タグの種類ごとに「ブラウザがどう表示するか」が決まっている",
        size=14.5, bold=True, color=INK, font=FONT_SANS)

footer(s, 16, total=21)

# ============================================================
# スライド 11.6: 発展トーク（サニタイズ・SQLインジェクション）
# ============================================================
s = add_slide()
set_bg(s)
kicker(s, "発展トーク")
title(s, "入力チェックのその先：サニタイズ")
textbox(s, Inches(0.5), Inches(1.35), Inches(11.5), Inches(0.5),
        "投稿の中身も、タグとしてそのまま実行できてしまったら？",
        size=15, color=SUB, font=FONT_SANS)

# ステップ1：もし入力がそのままタグとして実行されたら
card(s, Inches(0.5), Inches(2.05), Inches(12.3), Inches(1.9), fill=RGBColor(0xF7, 0xE9, 0xDC), line=MAIN)
textbox(s, Inches(0.85), Inches(2.25), Inches(11.5), Inches(0.45),
        "😱 投稿のタイトルに、さっきの <h1> をそのまま書けてしまったら？", size=16, bold=True, color=MAIN, font=FONT_SANS)
codebox = rect(s, Inches(0.85), Inches(2.8), Inches(11.6), Inches(0.65), fill=INK, radius=0.06)
tf = codebox.text_frame
tf.word_wrap = True
tf.margin_left = Inches(0.28); tf.margin_top = Inches(0.13)
p1 = tf.paragraphs[0]; p1.line_spacing = 1.2; p1.alignment = PP_ALIGN.LEFT
r1 = p1.add_run(); r1.text = "<h1>絶対見て！</h1>"
r1.font.name = "Courier New"; r1.font.size = Pt(15); r1.font.color.rgb = RGBColor(0xF7, 0xE9, 0xDC)
textbox(s, Inches(0.85), Inches(3.6), Inches(11.6), Inches(0.35),
        "→ 本物の見出しタグとして実行され、他の投稿と文字の大きさがバラバラになりレイアウトが崩れる",
        size=14, color=INK, font=FONT_SANS)

# ステップ2：対策（サニタイズ）
textbox(s, Inches(0.5), Inches(4.2), Inches(12.3), Inches(0.9),
        "この対策が「サニタイズ」：入力された文字をタグとして実行させず、\n無害な文字として表示に変換する処理（キミの旅の表示側はすでに対策済み）",
        size=14.5, bold=True, color=INK, font=FONT_SANS, line_spacing=1.4)

# ステップ3：SQLインジェクション（一言）
box = card(s, Inches(0.85), Inches(5.5), Inches(11.6), Inches(1.4), fill=SURFACE2)
textbox(s, Inches(1.15), Inches(5.68), Inches(11.0), Inches(1.05),
        "同じ考え方で、SQL文をそのまま実行できてしまうと、\nパスワードなどが盗まれる「SQLインジェクション」という攻撃もある。\n興味を持った人は、次回オープンキャンパスの「ホワイトハッカー体験」もチェック。",
        size=13.5, color=INK, font=FONT_SANS, line_spacing=1.4)

footer(s, 17, total=21)

# ============================================================
# スライド 12: アレンジ課題
# ============================================================
s = add_slide()
set_bg(s)
kicker(s, "STEP 5")
title(s, "自由にアレンジしてみよう")
textbox(s, Inches(0.5), Inches(1.32), Inches(7.5), Inches(0.4),
        "例：ページ全体の背景色を変えてみる",
        size=13.5, color=SUB, font=FONT_SANS)
code_chip(s, Inches(9.5), Inches(1.28), Inches(2.3), Inches(0.42), "style.css")

screenshot(s, os.path.join(IMG_DIR, "06_テーマカラー_修正前_オレンジ.png"), Inches(0.4), Inches(2.15),
           Inches(5.7), Inches(2.75), label="BEFORE：ベージュの背景", label_color=SUB)
arrow(s, Inches(6.15), Inches(3.35), Inches(0.85), Inches(0.5), fill=MAIN)
screenshot(s, os.path.join(IMG_DIR, "06_テーマカラー_修正後_ブルー.png"), Inches(7.2), Inches(2.15),
           Inches(5.7), Inches(2.75), label="AFTER：好きな色に変更", label_color=MAIN, border=MAIN)

code_block(s, Inches(2.5), Inches(5.15), Inches(8.3), Inches(0.9), [
    (13, "--main-color: #b5551f;", False),
    (14, "--main-color-hover: #96461a;", False),
    (16, "--bg: #f6f1e6;", True),
], code_size=11.5)

box = card(s, Inches(0.85), Inches(6.15), Inches(11.6), Inches(0.9), fill=RGBColor(0xF7, 0xE9, 0xDC), line=MAIN)
textbox(s, Inches(1.15), Inches(6.28), Inches(11.0), Inches(0.65),
        "ここまでは「Webプログラマー」の仕事。配色を工夫するのは「Webデザイナー」の仕事の入口。\n興味がある人は、次回オープンキャンパスの「Webデザイン体験」もチェックしてみよう。",
        size=12.5, color=INK, font=FONT_SANS, line_spacing=1.3)

footer(s, 18, total=21)

# ============================================================
# スライド 13: 体験の流れ（振り返り）
# ============================================================
s = add_slide()
set_bg(s)
kicker(s, "REVIEW")
title(s, "今日たどった道のり")

flow = [
    "プログラミングと Web の仕組みを知る",
    "サイトを起動して画面を確認",
    "ボタンの文字を直して“味見”",
    "会員登録してログイン",
    "現在地から住所を自動入力（API体験）",
    "コメント投稿を有効化",
    "データベースの考え方を知る",
    "入力チェックを自分で追加",
]
col_n = 2
row_n = 4
col_w = Inches(5.55)
row_h = Inches(1.0)
gap_x = Inches(0.2)
gap_y = Inches(0.12)
grid_w = Emu(col_w * col_n + gap_x)
start_x = Emu((SLIDE_W - grid_w) // 2)
start_y = Inches(1.85)

for i, f in enumerate(flow):
    col = i // row_n
    row = i % row_n
    x = Emu(start_x + col * (col_w + gap_x))
    y = Emu(start_y + row * (row_h + gap_y))
    card(s, x, y, col_w, row_h)
    step_badge(s, Emu(x + Inches(0.22)), Emu(y + Inches(0.29)), i + 1, "")
    textbox(s, Emu(x + Inches(0.85)), y, Emu(col_w - Inches(1.1)), row_h,
            f, size=13.5, color=INK, font=FONT_SANS, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.15)

footer(s, 19, total=21)

# ============================================================
# スライド 14: まとめ
# ============================================================
s = add_slide()
set_bg(s)
kicker(s, "SUMMARY")
title(s, "今日の体験、実は仕事そのもの")

pairs = [
    ("ボタンの文字を変える", "→ UIの改善"),
    ("住所を自動入力する", "→ 外部API連携・データ処理"),
    ("コメントを投稿できるようにする", "→ 機能の実装"),
    ("空欄チェックを追加する", "→ バリデーション（入力チェック）"),
    ("表示をアレンジする", "→ 使いやすさ・見た目の工夫"),
]
row_w = Inches(11.3)
row_x = Emu((SLIDE_W - row_w) // 2)
y = Inches(1.95)
for left, right in pairs:
    card(s, row_x, y, row_w, Inches(0.75))
    textbox(s, Emu(row_x + Inches(0.35)), y, Inches(6.3), Inches(0.75), left, size=14, color=INK,
            anchor=MSO_ANCHOR.MIDDLE, font=FONT_SANS)
    textbox(s, Emu(row_x + Inches(6.8)), y, Inches(4.3), Inches(0.75), right, size=14, color=MAIN, bold=True,
            anchor=MSO_ANCHOR.MIDDLE, font=FONT_SANS)
    y = Emu(y + Inches(0.85))

textbox(s, row_x, Inches(6.35), row_w, Inches(0.5),
        "小さな修正をして、動かして、確認する。この繰り返しが、実際の開発でも基本の流れです。",
        13.5, color=SUB, font=FONT_SANS, align=PP_ALIGN.CENTER)

footer(s, 20, total=21)

# ============================================================
# スライド 15: 締め
# ============================================================
s = add_slide()
set_bg(s, BG)
c1 = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(-2.2), Inches(4.5), Inches(5), Inches(5))
c1.fill.solid(); c1.fill.fore_color.rgb = RGBColor(0xEF, 0xE6, 0xD3)
c1.line.fill.background(); c1.shadow.inherit = False

c2 = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(10.5), Inches(-2), Inches(4.5), Inches(4.5))
c2.fill.solid(); c2.fill.fore_color.rgb = RGBColor(0xF7, 0xE9, 0xDC)
c2.line.fill.background(); c2.shadow.inherit = False

textbox(s, Inches(1.0), Inches(2.75), Inches(11.33), Inches(1.0),
        "小さな修正の積み重ねが、", size=28, color=INK, bold=True, align=PP_ALIGN.CENTER, font=FONT_SERIF)
textbox(s, Inches(1.0), Inches(3.4), Inches(11.33), Inches(1.0),
        "Webアプリを動かす力になる。", size=28, color=MAIN, bold=True, align=PP_ALIGN.CENTER, font=FONT_SERIF)

line = rect(s, Inches(5.87), Inches(4.5), Inches(1.6), Pt(3), fill=MAIN)

textbox(s, Inches(1.0), Inches(4.85), Inches(11.33), Inches(0.6),
        "ご参加ありがとうございました", size=16, color=SUB, align=PP_ALIGN.CENTER, font=FONT_SANS)

footer(s, 21, total=21)

out_path = "/Users/shoyabushita/Desktop/web_taiken/Webプログラミング体験_資料.pptx"
prs.save(out_path)
print("saved:", out_path, "slides:", len(prs.slides))
